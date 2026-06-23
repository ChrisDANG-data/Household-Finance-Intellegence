"""Generate Household Financial Intelligence presentation (.pptx)."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "Household_Financial_Intelligence.pptx"

SLIDES = [
    {
        "title": "Household Financial Intelligence",
        "subtitle": "AI-Powered Personal Finance — Final Project",
        "bullets": [
            "2026 Inference AI Course",
            "[Your name]",
            "One ledger · live Plaid · documents · mobile access",
        ],
        "notes": (
            "Introduce the problem: money is scattered across apps and PDFs; "
            "generic chatbots guess. This app centralizes truth and uses AI to explain."
        ),
    },
    {
        "title": "Problem",
        "subtitle": "Simple questions are hard to answer",
        "bullets": [
            "Money is everywhere — spreadsheets, banks, PDF bills",
            "Can we afford $20K travel in July? — no single simulation",
            "Contracts sit in folders — not in a ledger",
            "Need answers on the phone, not only the laptop",
        ],
        "notes": "Use the China trip and property tax examples from class.",
    },
    {
        "title": "Solution",
        "subtitle": "Four pillars",
        "bullets": [
            "01 PostgreSQL ledger: manual events + PDF obligations → FinancialEvent",
            "02 Plaid Direct API: sync cash + forecast; AI explains (engines compute)",
            "03 Documents → RAG for AI + Obsidian wiki for humans",
            "04 Telegram → n8n → Vercel API (ngrok for local dev only)",
        ],
        "notes": (
            "Correct Plaid: Direct API not MCP. ngrok only because n8n is localhost. "
            "Deterministic math first; LLM narrates."
        ),
    },
    {
        "title": "Four engines",
        "subtitle": "Strict boundaries",
        "bullets": [
            "Document Intelligence — upload, OCR, RAG",
            "Financial State — canonical ledger (Postgres)",
            "Forecast Simulation — deterministic projection.ts",
            "AI Explanation — read-only snapshots; never owns math",
        ],
        "notes": "See docs/ARCHITECTURE.md for dependency diagram.",
    },
    {
        "title": "Infrastructure",
        "subtitle": "Hybrid cloud + local automation",
        "bullets": [
            "Vercel — Next.js app (household-financial-web.vercel.app)",
            "Neon PostgreSQL + pgvector · Vercel Blob for files",
            "Plaid Direct API (sandbox) · LangGraph optional on Railway",
            "Local: Docker (n8n + Postgres) · Telegram via ngrok → n8n",
        ],
        "notes": "n8n does not auto-deploy; only Vercel on git push to main.",
    },
    {
        "title": "Demo",
        "subtitle": "Hard questions",
        "bullets": [
            "$20K China in July → what-if forecast + closing balance",
            "Property tax in June → ledger lookup + disposable on /balances",
            "Insurance bill terms → RAG on uploaded PDF",
            "Path: /ledger → /balances → /simulation (Ask AI)",
        ],
        "notes": "Have Plaid synced and ledger populated before demo.",
    },
    {
        "title": "Thank you",
        "subtitle": "Questions?",
        "bullets": [
            "Household Financial Intelligence",
            "Presentation/PRESENTATION_QA.md — 30 rehearsed answers",
            "Live: household-financial-web.vercel.app",
        ],
        "notes": "Offer to show Telegram if n8n is running.",
    },
]


def add_slide(prs: Presentation, slide_data: dict) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = slide_data["title"]

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    if slide_data.get("subtitle"):
        p = tf.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.level = 0
        p.font.size = Pt(18)

    for i, bullet in enumerate(slide_data.get("bullets", [])):
        p = tf.add_paragraph() if i > 0 or slide_data.get("subtitle") else tf.paragraphs[0]
        if i == 0 and not slide_data.get("subtitle"):
            p.text = bullet
        else:
            if i == 0 and slide_data.get("subtitle"):
                p = tf.add_paragraph()
            p.text = bullet
        p.level = 0
        p.font.size = Pt(16)

    if slide_data.get("notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        add_slide(prs, slide_data)

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
